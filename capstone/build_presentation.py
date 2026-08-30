"""
Build Assessment Task 3 Presentation — OENG1167
Cybersecurity for Connected Cars: Edge-Based Federated Intrusion Detection
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────
RMIT_RED    = RGBColor(0xE6, 0x00, 0x28)
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG2    = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x82, 0xC8)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xAB)
ACCENT_GREEN= RGBColor(0x2E, 0xCC, 0x71)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY    = RGBColor(0x99, 0x99, 0x99)
CARD_BG     = RGBColor(0x22, 0x2B, 0x45)
CARD_BG2    = RGBColor(0x1E, 0x25, 0x3C)
ORANGE      = RGBColor(0xF3, 0x96, 0x21)
YELLOW      = RGBColor(0xF1, 0xC4, 0x0F)
SOFT_RED    = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

FIGURES = "/Users/likhithgowda/Documents/capstone-cv2x-ids"

# ── Helper functions ────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=ACCENT_TEAL):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        # bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "▸  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"
        # text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tf

def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, width, color=RMIT_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    return False

def slide_number_footer(slide, num, total=14):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
             f"{num}/{total}", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def section_header(slide, section_text, color=ACCENT_TEAL):
    add_text(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
             section_text.upper(), font_size=11, color=color, bold=True)

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Red accent bar at top
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5),
         "Cybersecurity for Connected Cars", font_size=42, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.8),
         "Edge-Based Federated Intrusion Detection", font_size=28, color=ACCENT_TEAL)

add_accent_line(slide, Inches(0.8), Inches(4.0), Inches(3), RMIT_RED)

add_text(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
         "OENG1167 Engineering Capstone Project — Part A", font_size=16, color=LIGHT_GRAY)
add_text(slide, Inches(0.8), Inches(4.9), Inches(11), Inches(0.4),
         "Supervisor: A/Prof Ke (Desmond) Wang  |  Mentor: Mr. Kanwardeep Singh Gahlot",
         font_size=14, color=MID_GRAY)

add_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(1.2),
         "Likhith Lokesh Gowda  ·  Verna Nakhla  ·  Joshua Wong  ·  Ken Navarro  ·  Andrew Ng",
         font_size=16, color=LIGHT_GRAY)
add_text(slide, Inches(0.8), Inches(6.1), Inches(11), Inches(0.4),
         "School of Engineering — RMIT University", font_size=13, color=MID_GRAY)

slide_number_footer(slide, 1)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM & SIGNIFICANCE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Project Significance")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Why Connected Vehicle Security Matters", font_size=32, color=WHITE, bold=True)

# Left column — the problem
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8))
add_text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "The Problem", font_size=20, color=RMIT_RED, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(3.8), [
    "Connected vehicles broadcast Basic Safety Messages (BSMs) at 10 Hz",
    "BSMs carry position, speed, and heading — unauthenticated by design",
    "100ms latency constraint prevents heavyweight crypto",
    "Vulnerable to position spoofing, replay, false data injection, Sybil, DoS",
    "A single compromised BSM can trigger false emergency braking",
], font_size=15)

# Right column — why FL
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8))
add_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Why Federated Learning?", font_size=20, color=ACCENT_TEAL, bold=True)

add_bullet_list(slide, Inches(7.1), Inches(2.7), Inches(5.2), Inches(3.8), [
    "Centralised IDS requires streaming raw GPS data → privacy violation",
    "GDPR / Australian Privacy Act restrict GPS trajectory collection",
    "FL trains locally on each RSU, shares only model weights",
    "Raw BSM data never leaves the edge node",
    "Detection must run in <100ms on edge hardware (PC5 constraint)",
], font_size=15)

slide_number_footer(slide, 2)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — STAKEHOLDERS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Project Significance")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Stakeholders", font_size=32, color=WHITE, bold=True)

stakeholders = [
    ("Road Authorities", "Transport agencies (e.g. VicRoads, NHTSA) need V2X security standards", ACCENT_BLUE),
    ("Vehicle OEMs", "Manufacturers integrating C-V2X must protect BSM integrity", ACCENT_TEAL),
    ("Telecom Operators", "5G infrastructure providers deploying RSU edge compute", ORANGE),
    ("Road Users", "Drivers and pedestrians relying on safety-critical BSM data", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(stakeholders):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(2.0) + row * Inches(2.5)
    add_card(slide, x, y, Inches(5.7), Inches(2.0))
    # color bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(6), Inches(2.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.0), Inches(0.5),
             title, font_size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.0), Inches(1.0),
             desc, font_size=15, color=LIGHT_GRAY)

slide_number_footer(slide, 3)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — EXISTING TECHNOLOGIES / LITERATURE GAP
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Existing Technologies")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Literature Landscape & Research Gap", font_size=32, color=WHITE, bold=True)

# Existing approaches
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8))
add_text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Existing Approaches", font_size=20, color=ACCENT_BLUE, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(3.8), [
    "VeReMi dataset (2018) — VANET-only, no 5G NR, limited attack types",
    "MATLAB Classification Learner — 54.6% accuracy, dominant-class bias",
    "Centralised deep learning IDS — requires raw data collection",
    "Existing FL-IDS work uses WiFi/LTE, not C-V2X over 5G NR",
    "No published FL-IDS evaluates non-IID impact on C-V2X detection",
], font_size=15)

# Our contribution
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.8))
add_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Our Contribution", font_size=20, color=ACCENT_GREEN, bold=True)

add_bullet_list(slide, Inches(7.1), Inches(2.7), Inches(5.2), Inches(3.8), [
    "First 5G NR C-V2X IDS dataset (NS-3 + 5G-LENA, 12 attack scenarios)",
    "Dual-layer feature architecture: network + vehicular features",
    "Systematic non-IID evaluation: 60 FL experiments (Dirichlet + scenario)",
    "Edge deployment analysis proving <100ms feasibility",
    "Privacy-preserving: raw BSMs never leave the RSU",
], font_size=15, bullet_color=ACCENT_GREEN)

slide_number_footer(slide, 4)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — RESEARCH QUESTIONS & METHODOLOGY
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Project Design & Engineering Methods")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Research Questions & Methodology", font_size=32, color=WHITE, bold=True)

rqs = [
    ("RQ1", "Dataset Generation", "How to generate a labelled 5G C-V2X dataset\nthrough NS-3 simulation?", ACCENT_BLUE),
    ("RQ2", "Feature Engineering\n& Classification", "Which features best discriminate\nC-V2X attack types?", ACCENT_TEAL),
    ("RQ3", "Federated Learning", "How does non-IID data heterogeneity\naffect FedAvg convergence?", ORANGE),
    ("RQ4", "Edge Deployment", "Can inference meet the 100ms\nPC5 deadline on edge hardware?", ACCENT_GREEN),
]

for i, (rq, title, question, color) in enumerate(rqs):
    x = Inches(0.5) + i * Inches(3.15)
    add_card(slide, x, Inches(2.0), Inches(2.95), Inches(5.0))
    # RQ badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), Inches(2.3), Inches(0.8), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = rq
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"

    add_text(slide, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.8),
             title, font_size=16, color=color, bold=True)
    add_text(slide, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(1.5),
             question, font_size=13, color=LIGHT_GRAY)

slide_number_footer(slide, 5)

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — SYSTEM ARCHITECTURE PIPELINE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Project Design & Engineering Methods")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Four-Stage Engineering Pipeline", font_size=32, color=WHITE, bold=True)

stages = [
    ("1", "Dataset\nGeneration", "NS-3 + 5G-LENA\n40 UEs, 4 gNBs\n12 scenarios\n18,240 rows", ACCENT_BLUE),
    ("2", "Feature\nEngineering", "ANOVA + MI ranking\nBorda count fusion\n39 → 15 features\nSHAP validation", ACCENT_TEAL),
    ("3", "Centralised\nClassification", "RF / GBC / MLP\nAll F1 = 1.00\nModel spec export\nfor FL contract", ORANGE),
    ("4", "Federated\nLearning", "FedAvg, 60 exps\nDirichlet + scenario\nLatency: 26.4 μs\n3,789× headroom", ACCENT_GREEN),
]

for i, (num, title, details, color) in enumerate(stages):
    x = Inches(0.5) + i * Inches(3.15)
    add_card(slide, x, Inches(2.0), Inches(2.85), Inches(5.0))

    # stage number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), Inches(2.3), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"

    add_text(slide, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(0.8),
             title, font_size=17, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(4.0), Inches(2.5), Inches(2.5),
             details, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # arrow between stages
    if i < 3:
        arrow_x = x + Inches(2.95)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(4.0), Inches(0.2), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

slide_number_footer(slide, 6)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — TECHNICAL: DATASET GENERATION
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Technical Progress — RQ1")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "5G C-V2X Dataset Generation", font_size=32, color=WHITE, bold=True)

# Left — key specs
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0))
add_text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Simulation Specifications", font_size=18, color=ACCENT_BLUE, bold=True)

specs = [
    "NS-3.42 + 5G-LENA NR module (Band n78, 3.5 GHz, 20 MHz TDD)",
    "40 vehicles, 4 gNBs, BSMs at 10 Hz (ETSI CAM standard)",
    "12 scenarios: 1 benign + 5 network attacks + 6 vehicular attacks",
    "600s per scenario, per-scenario random seeds (43–54)",
    "30-second time-windowed feature aggregation",
    "18,240 rows × 39 columns — zero NaN/Inf",
    "70/15/15 stratified split with group leakage prevention",
    "57 internal integrity checks — all pass",
]
add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(4.0), specs, font_size=14)

# Right — label distribution figure
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0))
add_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Label Distribution", font_size=18, color=ACCENT_BLUE, bold=True)
add_image_safe(slide, f"{FIGURES}/dataset-expansion/output/figures/label_distribution.png",
               Inches(7.1), Inches(2.8), width=Inches(5.2))

slide_number_footer(slide, 7)

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNICAL: FEATURE ENGINEERING
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Technical Progress — RQ2")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Feature Engineering: 39 → 15 Features", font_size=32, color=WHITE, bold=True)

# Left — pipeline
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0))
add_text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Selection Pipeline", font_size=18, color=ACCENT_TEAL, bold=True)

pipeline_steps = [
    "39 raw columns → remove metadata, labels, zero-variance, context",
    "24 informative → correlation filter (ρ > 0.99) → 17",
    "ANOVA F-scores + Mutual Information + Borda count ranking",
    "Top-k sweep: k=15 achieves perfect F1; k=13 drops to 0.51",
    "Dual-layer: network features (flood_ratio, n_flood, pkt_rate)",
    "Dual-layer: vehicular features (pos_deviation, speed_deviation)",
    "Each attack maps to specific features (validated via SHAP)",
]
add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(4.0), pipeline_steps, font_size=14)

# Right — SHAP or top-k figure
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0))
add_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Top-k Feature Sweep", font_size=18, color=ACCENT_TEAL, bold=True)
add_image_safe(slide, f"{FIGURES}/feature-engineering/output/figures/topk_curves.png",
               Inches(7.1), Inches(2.8), width=Inches(5.2))

slide_number_footer(slide, 8)

# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNICAL: CLASSIFICATION
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Technical Progress — RQ2")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Centralised Classification Baseline", font_size=32, color=WHITE, bold=True)

# Results summary
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0))
add_text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Three-Model Comparison", font_size=18, color=ORANGE, bold=True)

class_points = [
    "Random Forest: Macro F1 = 1.0000, MCC = 1.0000",
    "Gradient Boosting: Macro F1 = 1.0000, MCC = 1.0000",
    "MLP [15→128→64→32→12]: Macro F1 = 1.0000, MCC = 1.0000",
    "All 12 classes at per-class F1 = 1.0, FPR = 0.0",
    "Perfect scores are legitimate: deterministic NS-3 attack signatures",
    "MLP architecture + scaler params exported for FL contract",
    "Dual-layer ablation: network-only F1 = 0.84, vehicular-only = 0.83",
    "Combined dual-layer F1 = 1.00 — proves both domains required",
]
add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(4.0), class_points, font_size=14)

# Right — confusion matrix
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0))
add_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "MLP Confusion Matrix (12-class)", font_size=18, color=ORANGE, bold=True)
add_image_safe(slide, f"{FIGURES}/classification/output/figures/confusion_mlp.png",
               Inches(7.1), Inches(2.8), width=Inches(5.2))

slide_number_footer(slide, 9)

# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECHNICAL: FEDERATED LEARNING RESULTS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Technical Progress — RQ3")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Federated Learning: 60 Experiments", font_size=32, color=WHITE, bold=True)

# Left — key findings
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0))
add_text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Key Findings", font_size=18, color=ACCENT_GREEN, bold=True)

fl_points = [
    "FedAvg with custom MLP, 50 global rounds, class-weighted loss",
    "IID (α=100): F1 = 1.00 — matches centralised baseline",
    "Mild non-IID (α ≥ 0.5): F1 = 1.00 — FedAvg is robust",
    "Strong non-IID (α=0.1, C=5, E=1): F1 = 0.54 — significant drop",
    "Extra local epochs (E=3) partially recover: F1 = 0.94",
    "Scenario-based: F1 = 0.34–0.56 across all configs",
    "More local epochs worsen scenario-based (clients overfit)",
    "FedAvg cannot overcome complete class gaps → FedProx needed",
]
add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(4.0), fl_points, font_size=14)

# Right — non-IID degradation figure
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0))
add_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Non-IID Degradation", font_size=18, color=ACCENT_GREEN, bold=True)
add_image_safe(slide, f"{FIGURES}/federated-learning/output/figures/noniid_degradation.png",
               Inches(7.1), Inches(2.8), width=Inches(5.2))

slide_number_footer(slide, 10)

# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — TECHNICAL: EDGE DEPLOYMENT
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Technical Progress — RQ4")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Edge Deployment Feasibility", font_size=32, color=WHITE, bold=True)

# Three metric cards
metrics = [
    ("26.4 μs", "Mean Inference", "2,000 forward passes\nSingle-core CPU", ACCENT_BLUE),
    ("3,789×", "Latency Headroom", "Under 100ms PC5\nconstraint", ACCENT_GREEN),
    ("908 params", "Min Viable Model", "[15→32→12] at 3.6 KB\nStill achieves F1 = 1.00", ORANGE),
]

for i, (big_num, label, detail, color) in enumerate(metrics):
    x = Inches(0.8) + i * Inches(4.1)
    add_card(slide, x, Inches(2.0), Inches(3.7), Inches(2.5))
    add_text(slide, x + Inches(0.2), Inches(2.2), Inches(3.3), Inches(0.8),
             big_num, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.0), Inches(3.3), Inches(0.4),
             label, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.5), Inches(3.3), Inches(0.8),
             detail, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Communication cost card
add_card(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0))
add_text(slide, Inches(1.1), Inches(5.1), Inches(5.0), Inches(0.5),
         "Communication Break-Even", font_size=18, color=ACCENT_TEAL, bold=True)
comm_points = [
    "FL model upload: 51 KB  vs  1 vehicle streaming raw BSMs for 17 seconds",
    "8 vehicles per RSU surpass total FL training cost in ~3.5 minutes",
    "FL never transmits raw BSM data — preserves driver location privacy (GDPR compliant)",
]
add_bullet_list(slide, Inches(1.1), Inches(5.6), Inches(11.0), Inches(1.2), comm_points, font_size=14)

slide_number_footer(slide, 11)

# ════════════════════════════════════════════════════════════════════
# SLIDE 12 — SUMMARY OF DELIVERABLES
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Technical Progress")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Part A Deliverables — All Complete", font_size=32, color=WHITE, bold=True)

deliverables = [
    ("RQ1", "Dataset Generation", "18,240 rows, 12 scenarios, 57 integrity checks pass", ACCENT_BLUE, "✓"),
    ("RQ2a", "Feature Engineering", "39 → 15 features, multiclass F1 = 1.00", ACCENT_TEAL, "✓"),
    ("RQ2b", "Classification", "3 models at F1 = 1.00, MLP spec exported for FL", ORANGE, "✓"),
    ("RQ3", "Federated Learning", "60 experiments, non-IID degradation characterised", ACCENT_GREEN, "✓"),
    ("RQ4", "Edge Deployment", "26.4 μs inference, 3,789× headroom", RMIT_RED, "✓"),
]

for i, (rq, title, metric, color, status) in enumerate(deliverables):
    y = Inches(2.0) + i * Inches(1.05)
    add_card(slide, Inches(0.8), y, Inches(11.7), Inches(0.9))

    # RQ badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), y + Inches(0.15), Inches(0.85), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = rq
    run.font.size = Pt(13)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"

    add_text(slide, Inches(2.2), y + Inches(0.1), Inches(2.5), Inches(0.5),
             title, font_size=17, color=WHITE, bold=True)
    add_text(slide, Inches(4.8), y + Inches(0.1), Inches(6.0), Inches(0.5),
             metric, font_size=15, color=LIGHT_GRAY)

    # checkmark
    add_text(slide, Inches(11.5), y + Inches(0.05), Inches(0.8), Inches(0.6),
             status, font_size=24, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 12)

# ════════════════════════════════════════════════════════════════════
# SLIDE 13 — TIMELINE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)
section_header(slide, "Timeline")

add_text(slide, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8),
         "Project Timeline: Part A → Part B", font_size=32, color=WHITE, bold=True)

# Part A timeline (completed)
add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0))
add_text(slide, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Part A — Completed", font_size=20, color=ACCENT_GREEN, bold=True)

part_a = [
    "Weeks 1–3: Literature review, project proposal, risk assessment",
    "Weeks 4–6: NS-3 simulation, dataset generation pipeline",
    "Weeks 7–9: Feature engineering, classification baselines",
    "Weeks 10–12: Federated learning experiments (60 configs)",
    "Week 12: Edge deployment analysis, report integration",
    "Week 13: Presentation delivery",
]
add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.0), Inches(4.0), part_a, font_size=14,
                bullet_color=ACCENT_GREEN)

# Part B priorities
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0))
add_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.5),
         "Part B — Planned", font_size=20, color=ORANGE, bold=True)

part_b = [
    "FedProx: proximal term to handle non-IID degradation",
    "Head-to-head FedAvg vs FedProx across all 20 configurations",
    "Krum aggregation: Byzantine-resilient alternative",
    "Statistical rigour: ≥6 seeds (Part A used 3)",
    "FPGA profiling on actual edge hardware (if available)",
    "Final report and Part B presentation",
]
add_bullet_list(slide, Inches(7.1), Inches(2.7), Inches(5.2), Inches(4.0), part_b, font_size=14,
                bullet_color=ORANGE)

slide_number_footer(slide, 13)

# ════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING / QUESTIONS
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, RMIT_RED)

add_text(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
         "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.8),
         "Questions?", font_size=32, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(4.5), Inches(2.3), RMIT_RED)

add_text(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
         "Cybersecurity for Connected Cars: Edge-Based Federated Intrusion Detection",
         font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.4),
         "OENG1167 Engineering Capstone Project — RMIT University",
         font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

slide_number_footer(slide, 14)

# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
out_dir = "/Users/likhithgowda/Library/Mobile Documents/com~apple~CloudDocs/Academics/Engineering Capstone Project/Assessment Task 3"
out_path = os.path.join(out_dir, "OENG1167_Assessment3_Presentation.pptx")
prs.save(out_path)
print(f"Saved to: {out_path}")
print(f"Slides: {len(prs.slides)}")
